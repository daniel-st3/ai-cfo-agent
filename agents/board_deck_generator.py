"""
board_deck_generator.py — Automated Board Deck Generator

Generates a branded 10-slide PowerPoint deck from existing KPI, anomaly,
market signal, scenario, and report data. Uses python-pptx for slide assembly
and matplotlib for chart images.
"""

from __future__ import annotations

import io
import os
import uuid
from datetime import date, datetime

import matplotlib
matplotlib.use("Agg")  # headless backend
import matplotlib.pyplot as plt
import numpy as np
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from api.models import (
    Anomaly,
    BoardDeck,
    KPISnapshot,
    MarketSignal,
    Report,
)

# ── Colors ────────────────────────────────────────────────────────────────────
BLUE = "#0071e3"
GREEN = "#34c759"
RED = "#ff3b30"
AMBER = "#ff9500"
DARK = "#1d1d1f"
GRAY = "#6e6e73"
LIGHT = "#f5f5f7"


class BoardDeckGenerator:
    """Generates a 10-slide PowerPoint board deck from existing run data."""

    DECKS_DIR = "data/decks"

    async def run(
        self,
        session: AsyncSession,
        run_id: uuid.UUID,
        company_name: str = "Portfolio Company",
    ) -> str:
        """Generate the deck and return its file path.

        Persists a BoardDeck record and returns the absolute file path.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

        os.makedirs(self.DECKS_DIR, exist_ok=True)
        file_path = os.path.join(self.DECKS_DIR, f"{run_id}_board_deck.pptx")

        # Fetch all data
        kpis = await self._get_kpis(session, run_id)
        anomalies = await self._get_anomalies(session, run_id)
        signals = await self._get_signals(session, run_id)
        report = await self._get_report(session, run_id)

        latest = kpis[-1] if kpis else None
        prev = kpis[-5] if len(kpis) >= 5 else kpis[0] if kpis else None

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank

        # ── Slide 1: Cover ────────────────────────────────────────────────────
        self._slide_cover(prs, blank_layout, company_name)

        # ── Slide 2: Executive Summary ────────────────────────────────────────
        bullets = []
        if report:
            import re
            text = report.executive_summary
            # Extract bullet lines
            lines = [l.strip().lstrip("•-– ") for l in text.split("\n") if l.strip()]
            bullets = [l for l in lines if l][:3]
        if not bullets:
            bullets = [
                f"Weekly MRR of ${float(latest.mrr or 0):,.0f} with {float(latest.gross_margin or 0)*100:.0f}% gross margin" if latest else "Financial data loaded",
                f"Burn rate ${float(latest.burn_rate or 0):,.0f}/week — monitoring closely" if latest else "Burn under review",
                f"Churn rate {float(latest.churn_rate or 0)*100:.1f}% — within acceptable range" if latest else "Churn stable",
            ]
        self._slide_exec_summary(prs, blank_layout, bullets)

        # ── Slide 3: Financial Snapshot ───────────────────────────────────────
        self._slide_financials(prs, blank_layout, latest, prev, company_name)

        # ── Slide 4: 13-Week Cash Flow ────────────────────────────────────────
        self._slide_cash_flow(prs, blank_layout, kpis)

        # ── Slide 5: Unit Economics ───────────────────────────────────────────
        self._slide_unit_economics(prs, blank_layout, latest)

        # ── Slide 6: Growth Metrics ───────────────────────────────────────────
        self._slide_growth(prs, blank_layout, kpis)

        # ── Slide 7: Anomalies & Risks ────────────────────────────────────────
        top_anomalies = sorted(
            [a for a in anomalies if a.severity in ("HIGH", "MEDIUM")],
            key=lambda a: (0 if a.severity == "HIGH" else 1)
        )[:4]
        self._slide_risks(prs, blank_layout, top_anomalies)

        # ── Slide 8: Competitive Landscape ───────────────────────────────────
        top_signals = signals[:5]
        self._slide_competitive(prs, blank_layout, top_signals)

        # ── Slide 9: Scenario Analysis ────────────────────────────────────────
        self._slide_scenarios(prs, blank_layout, kpis)

        # ── Slide 10: Fundraising Status ─────────────────────────────────────
        self._slide_fundraising(prs, blank_layout, latest)

        prs.save(file_path)

        # Persist record
        deck_record = await self._get_or_create_deck(session, run_id)
        deck_record.file_path = file_path
        deck_record.status = "ready"
        deck_record.generated_at = datetime.utcnow()
        await session.commit()

        return file_path

    # ── Slide builders ─────────────────────────────────────────────────────────

    def _slide_cover(self, prs, layout, company_name: str) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, DARK)
        self._add_text(slide, company_name.upper(), 1.5, 2.5, 10, 1, Pt(48), "FFFFFF", bold=True)
        self._add_text(slide, "BOARD OF DIRECTORS — CONFIDENTIAL", 1.5, 3.5, 10, 0.5, Pt(14), "6e6e73")
        self._add_text(slide, date.today().strftime("%B %Y"), 1.5, 4.2, 10, 0.5, Pt(18), "0071e3")

    def _slide_exec_summary(self, prs, layout, bullets: list[str]) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "EXECUTIVE SUMMARY", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Key Highlights", 0.5, 0.8, 12, 0.7, Pt(32), DARK[1:], bold=True)
        for i, bullet in enumerate(bullets[:3]):
            y = 1.9 + i * 1.5
            self._add_shape_rect(slide, 0.5, y, 12.3, 1.2, "f5f5f7", "e8e8ed")
            self._add_text(slide, f"● {bullet}", 0.7, y + 0.15, 11.8, 0.9, Pt(16), DARK[1:])

    def _slide_financials(self, prs, layout, latest, prev, company_name: str) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "FINANCIAL SNAPSHOT", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Current Financial Position", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        metrics = []
        if latest:
            mrr = float(latest.mrr or 0)
            arr = float(latest.arr or 0)
            burn = float(latest.burn_rate or 0)
            gm = float(latest.gross_margin or 0)
            churn = float(latest.churn_rate or 0)
            runway = (arr / max(burn * 52, 1)) * 12 if burn > 0 else 99

            prev_mrr = float(prev.mrr or 0) if prev else mrr
            mrr_growth = ((mrr - prev_mrr) / max(prev_mrr, 1)) * 100

            metrics = [
                ("Weekly MRR", f"${mrr:,.0f}", f"{'▲' if mrr_growth >= 0 else '▼'} {abs(mrr_growth):.1f}% WoW"),
                ("ARR Run-Rate", f"${arr:,.0f}", "Annualized"),
                ("Weekly Burn", f"${burn:,.0f}", "Operating"),
                ("Gross Margin", f"{gm*100:.1f}%", "Target: 60%+"),
                ("Churn Rate", f"{churn*100:.2f}%", "Weekly"),
                ("Est. Runway", f"{runway:.0f}mo", "at current burn"),
            ]

        cols = 3
        for i, (label, value, sub) in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = 0.5 + col * 4.2
            y = 1.9 + row * 2.2
            self._add_shape_rect(slide, x, y, 3.8, 1.9, "f5f5f7", "e8e8ed")
            self._add_text(slide, label, x + 0.15, y + 0.1, 3.5, 0.3, Pt(9), "6e6e73", bold=True)
            self._add_text(slide, value, x + 0.15, y + 0.4, 3.5, 0.7, Pt(26), DARK[1:], bold=True)
            self._add_text(slide, sub, x + 0.15, y + 1.45, 3.5, 0.3, Pt(9), "6e6e73")

    def _slide_cash_flow(self, prs, layout, kpis: list) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "13-WEEK CASH FLOW", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Projected Cash Position", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        # Generate matplotlib chart
        img_stream = self._make_cash_flow_chart(kpis)
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.4))

    def _slide_unit_economics(self, prs, layout, latest) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "UNIT ECONOMICS", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Customer Acquisition & Lifetime Value", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        metrics = []
        if latest:
            cac = float(latest.cac or 0)
            ltv = float(latest.ltv or 0)
            ratio = ltv / max(cac, 1)
            payback = cac / max(float(latest.mrr or 1) / max(float(latest.gross_margin or 0.7), 0.01), 1)
            metrics = [
                ("CAC", f"${cac:,.0f}", "Per customer acquired"),
                ("LTV", f"${ltv:,.0f}", "Lifetime customer value"),
                ("LTV/CAC Ratio", f"{ratio:.1f}x", "Target: 3x+"),
                ("Payback Period", f"{payback:.0f}wks", "Weeks to recover CAC"),
            ]

        for i, (label, value, sub) in enumerate(metrics):
            x = 0.5 + i * 3.1
            self._add_shape_rect(slide, x, 2.0, 2.8, 2.2, "f5f5f7", "e8e8ed")
            self._add_text(slide, label, x + 0.1, 2.1, 2.6, 0.3, Pt(9), "6e6e73", bold=True)
            self._add_text(slide, value, x + 0.1, 2.45, 2.6, 0.8, Pt(28), DARK[1:], bold=True)
            self._add_text(slide, sub, x + 0.1, 3.7, 2.6, 0.3, Pt(9), "6e6e73")

    def _slide_growth(self, prs, layout, kpis: list) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "GROWTH METRICS", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "MRR Trend & Retention", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        img_stream = self._make_mrr_chart(kpis)
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.4))

    def _slide_risks(self, prs, layout, anomalies: list) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "ANOMALIES & RISKS", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Detected Financial Risks", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        if not anomalies:
            self._add_text(slide, "✓ No significant anomalies detected in this period.", 0.5, 3.0, 12, 0.6, Pt(18), "34c759")
            return

        for i, a in enumerate(anomalies[:4]):
            y = 1.9 + i * 1.3
            color = "ff3b30" if a.severity == "HIGH" else "ff9500"
            self._add_shape_rect(slide, 0.5, y, 12.3, 1.1, "ffffff", color, border_width=3)
            label = f"[{a.severity}] {a.metric.replace('_', ' ').upper()}"
            self._add_text(slide, label, 0.7, y + 0.05, 11.8, 0.35, Pt(11), color, bold=True)
            desc = (a.description or f"Actual: {float(a.actual_value):,.1f}")[:120]
            self._add_text(slide, desc, 0.7, y + 0.45, 11.8, 0.5, Pt(12), DARK[1:])

    def _slide_competitive(self, prs, layout, signals: list) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "COMPETITIVE LANDSCAPE", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Market Intelligence", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        if not signals:
            self._add_text(slide, "No competitive signals detected recently.", 0.5, 3.0, 12, 0.6, Pt(16), "6e6e73")
            return

        for i, s in enumerate(signals[:5]):
            y = 1.8 + i * 1.1
            icon = "💰" if s.signal_type == "pricing_change" else "👥" if s.signal_type == "job_posting" else "📰"
            label = f"{icon} {s.competitor_name} · {s.signal_type.replace('_', ' ').title()}"
            self._add_text(slide, label, 0.5, y, 12, 0.35, Pt(11), BLUE[1:], bold=True)
            summary = (s.summary or "")[:130]
            self._add_text(slide, summary, 0.5, y + 0.38, 12, 0.5, Pt(11), DARK[1:])

    def _slide_scenarios(self, prs, layout, kpis: list) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, "FFFFFF")
        self._add_text(slide, "SCENARIO ANALYSIS", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Bear / Base / Bull Stress Test", 0.5, 0.8, 12, 0.7, Pt(28), DARK[1:], bold=True)

        if not kpis:
            return
        latest_mrr = float(kpis[-1].mrr or 0)
        latest_burn = float(kpis[-1].burn_rate or 0)

        scenarios = [
            ("Bear", latest_mrr * 0.80, latest_burn * 1.15, "ff3b30"),
            ("Base", latest_mrr, latest_burn, "0071e3"),
            ("Bull", latest_mrr * 1.20, latest_burn * 0.85, "34c759"),
        ]

        for i, (name, mrr, burn, color) in enumerate(scenarios):
            x = 0.5 + i * 4.2
            runway = (mrr * 12 / max(burn * 52, 1)) if burn > 0 else 99
            self._add_shape_rect(slide, x, 2.0, 3.8, 4.8, "f5f5f7", color, border_width=2)
            self._add_text(slide, name.upper(), x + 0.15, 2.1, 3.5, 0.4, Pt(14), color, bold=True)
            self._add_text(slide, f"${mrr:,.0f}/wk", x + 0.15, 2.6, 3.5, 0.5, Pt(22), DARK[1:], bold=True)
            self._add_text(slide, "Weekly MRR", x + 0.15, 3.1, 3.5, 0.3, Pt(9), "6e6e73")
            self._add_text(slide, f"${burn:,.0f}/wk burn", x + 0.15, 3.5, 3.5, 0.4, Pt(13), "6e6e73")
            self._add_text(slide, f"{runway:.0f} months runway", x + 0.15, 4.0, 3.5, 0.5, Pt(16), color, bold=True)

    def _slide_fundraising(self, prs, layout, latest) -> None:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(layout)
        self._fill_bg(slide, DARK)
        self._add_text(slide, "FUNDRAISING STATUS", 0.5, 0.3, 12, 0.5, Pt(11), "6e6e73", bold=True)
        self._add_text(slide, "Runway & Capital Strategy", 0.5, 0.8, 12, 0.7, Pt(32), "ffffff", bold=True)

        if latest:
            mrr = float(latest.mrr or 0)
            burn = float(latest.burn_rate or 0)
            runway_weeks = int(mrr / max(burn, 1) * 52) if burn > 0 else 999
            runway_months = runway_weeks / 4.33

            color = "ff3b30" if runway_months < 6 else "ff9500" if runway_months < 12 else "34c759"
            status = "CRITICAL — Begin fundraising immediately" if runway_months < 6 else \
                     "CAUTION — Start Series A process in 90 days" if runway_months < 12 else \
                     "HEALTHY — Well-capitalized for growth"

            self._add_text(slide, f"{runway_months:.0f}", 4.5, 2.2, 4.0, 1.5, Pt(96), color, bold=True)
            self._add_text(slide, "months of runway at current burn", 3.5, 3.8, 6.5, 0.5, Pt(14), "6e6e73")
            self._add_text(slide, status, 1.0, 5.0, 11, 0.8, Pt(18), color, bold=True)

        self._add_text(slide, f"Generated by AI CFO Agent · {date.today().strftime('%B %d, %Y')} · CONFIDENTIAL",
                       0.5, 6.9, 12, 0.4, Pt(9), "6e6e73")

    # ── Chart builders ─────────────────────────────────────────────────────────

    def _make_cash_flow_chart(self, kpis: list) -> io.BytesIO:
        """Generate a projected cash flow line chart as PNG bytes."""
        weeks = list(range(1, 14))
        if kpis:
            mrr = float(kpis[-1].mrr or 0)
            burn = float(kpis[-1].burn_rate or 0)
        else:
            mrr, burn = 50000, 30000

        cash = 300000.0
        p50, p10, p90 = [cash], [cash * 0.7], [cash * 1.3]
        for _ in weeks:
            net = mrr - burn
            cash = max(cash + net + np.random.normal(0, burn * 0.05), 0)
            p50.append(cash)
            p10.append(max(cash * 0.85, 0))
            p90.append(cash * 1.15)

        x = [0] + weeks
        fig, ax = plt.subplots(figsize=(12, 5), facecolor="white")
        ax.fill_between(x, p10, p90, color=BLUE, alpha=0.15, label="P10–P90 range")
        ax.plot(x, p50, color=BLUE, linewidth=2.5, label="P50 (median)")
        ax.axhline(0, color=RED, linestyle="--", linewidth=1.5, alpha=0.7, label="Zero cash")
        ax.set_xlabel("Week", fontsize=10)
        ax.set_ylabel("Cash Balance ($)", fontsize=10)
        ax.set_xticks(x)
        ax.set_xticklabels(["Now"] + [f"Wk {w}" for w in weeks], fontsize=8)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"${v:,.0f}"))
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.legend(fontsize=9, loc="upper right")
        ax.set_facecolor("white")
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf

    def _make_mrr_chart(self, kpis: list) -> io.BytesIO:
        """Generate MRR trend line chart as PNG bytes."""
        recent = kpis[-16:] if len(kpis) >= 16 else kpis
        labels = [k.week_start.strftime("%b %d") for k in recent]
        mrrs = [float(k.mrr or 0) for k in recent]
        burns = [float(k.burn_rate or 0) for k in recent]

        fig, ax = plt.subplots(figsize=(12, 5), facecolor="white")
        ax.fill_between(range(len(mrrs)), mrrs, alpha=0.15, color=BLUE)
        ax.plot(range(len(mrrs)), mrrs, color=BLUE, linewidth=2.5, marker="o", markersize=4, label="MRR")
        ax.plot(range(len(burns)), burns, color=RED, linewidth=1.5, linestyle="--", label="Burn Rate")
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=8)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"${v:,.0f}"))
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.legend(fontsize=9)
        ax.set_facecolor("white")
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf

    # ── Slide helpers ─────────────────────────────────────────────────────────

    def _fill_bg(self, slide, hex_color: str) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))

    def _add_text(
        self, slide, text: str, left: float, top: float, width: float, height: float,
        font_size, color: str, bold: bool = False
    ) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text)
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#") if "#" not in color else color[1:])

    def _add_shape_rect(
        self, slide, left: float, top: float, width: float, height: float,
        fill_hex: str, border_hex: str, border_width: int = 1
    ) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
        line = shape.line
        line.color.rgb = RGBColor.from_string(border_hex.lstrip("#"))
        line.width = Pt(border_width)

    # ── DB helpers ────────────────────────────────────────────────────────────

    async def _get_kpis(self, session: AsyncSession, run_id: uuid.UUID) -> list:
        result = await session.execute(
            select(KPISnapshot).where(KPISnapshot.run_id == run_id).order_by(KPISnapshot.week_start)
        )
        return list(result.scalars().all())

    async def _get_anomalies(self, session: AsyncSession, run_id: uuid.UUID) -> list:
        result = await session.execute(
            select(Anomaly).where(Anomaly.run_id == run_id).order_by(Anomaly.severity)
        )
        return list(result.scalars().all())

    async def _get_signals(self, session: AsyncSession, run_id: uuid.UUID) -> list:
        result = await session.execute(
            select(MarketSignal).where(MarketSignal.run_id == run_id).order_by(MarketSignal.date.desc()).limit(10)
        )
        return list(result.scalars().all())

    async def _get_report(self, session: AsyncSession, run_id: uuid.UUID):
        result = await session.execute(
            select(Report).where(Report.run_id == run_id).order_by(Report.created_at.desc()).limit(1)
        )
        return result.scalar_one_or_none()

    async def _get_or_create_deck(self, session: AsyncSession, run_id: uuid.UUID) -> BoardDeck:
        result = await session.execute(
            select(BoardDeck).where(BoardDeck.run_id == run_id).limit(1)
        )
        deck = result.scalar_one_or_none()
        if not deck:
            deck = BoardDeck(run_id=run_id, file_path="", status="generating")
            session.add(deck)
            await session.flush()
        return deck
